"""
TrustOffice Branded PowerPoint Template Generator
Creates 5 slide templates: Title, Quote, List, Image Left, Image Right
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x01, 0x00, 0x79)
GOLD = RGBColor(0xD5, 0xAD, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xF8, 0xF7, 0xF4)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
GOLD_LIGHT = RGBColor(0xF5, 0xEE, 0xD5)

W = Inches(13.333)
H = Inches(7.5)

FONT_SERIF = "Georgia"
FONT_SANS = "Helvetica Neue"
FONT_MONO = "Courier New"


def add_navy_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_light_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_gold_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_footer(slide, dark=False):
    """Add TrustOffice footer bar"""
    # Footer background
    footer_top = H - Inches(0.55)
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, footer_top, W, Inches(0.55))
    if dark:
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x50)
    else:
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    footer.line.fill.background()

    # Gold top line on footer
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, footer_top, W, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Footer text
    tf = slide.shapes.add_textbox(Inches(0.8), footer_top + Inches(0.12), Inches(5), Inches(0.3))
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "TrustOffice"
    run.font.name = FONT_SERIF
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = GOLD if dark else NAVY

    run2 = p.add_run()
    run2.text = "  |  Trust Governance Workspace"
    run2.font.name = FONT_SANS
    run2.font.size = Pt(9)
    run2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xCC) if dark else GRAY_TEXT

    # Right side
    tf2 = slide.shapes.add_textbox(W - Inches(3.5), footer_top + Inches(0.12), Inches(2.7), Inches(0.3))
    p2 = tf2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run3 = p2.add_run()
    run3.text = "trustoffice.app"
    run3.font.name = FONT_MONO
    run3.font.size = Pt(9)
    run3.font.color.rgb = GOLD if dark else NAVY


def add_corner_accent(slide, dark=False):
    """Add a small gold square accent in top-left"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.6), Inches(0.35), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def slide_1_title(prs):
    """TITLE SLIDE — Navy background, centered, bold"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_navy_bg(slide)

    # Top gold line
    add_gold_accent_line(slide, 0, 0, W)

    # Side gold bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.2), Inches(0.08), Inches(2.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    # Title
    tf = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10), Inches(1.5))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Presentation Title"
    run.font.name = FONT_SERIF
    run.font.size = Pt(48)
    run.font.color.rgb = WHITE
    run.font.bold = False

    # Subtitle
    tf2 = slide.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(10), Inches(0.8))
    tf2.text_frame.word_wrap = True
    p2 = tf2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Subtitle or description goes here"
    run2.font.name = FONT_SANS
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xDD)

    # Bottom gold line
    add_gold_accent_line(slide, Inches(1.2), Inches(5.0), Inches(3))

    # Logo / branding area bottom-left
    tf3 = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(4), Inches(0.5))
    p3 = tf3.text_frame.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "TrustOffice"
    run3.font.name = FONT_SERIF
    run3.font.size = Pt(16)
    run3.font.bold = True
    run3.font.color.rgb = GOLD

    run4 = p3.add_run()
    run4.text = "  |  Date"
    run4.font.name = FONT_SANS
    run4.font.size = Pt(12)
    run4.font.color.rgb = RGBColor(0x88, 0x88, 0xBB)


def slide_2_quote(prs):
    """QUOTE SLIDE — Light background, large centered quote"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)

    # Top navy bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Large open-quote mark
    tf_q = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(2), Inches(1.5))
    p_q = tf_q.text_frame.paragraphs[0]
    run_q = p_q.add_run()
    run_q.text = "\u201C"
    run_q.font.name = FONT_SERIF
    run_q.font.size = Pt(120)
    run_q.font.color.rgb = GOLD

    # Quote text
    tf = slide.shapes.add_textbox(Inches(2.0), Inches(2.4), Inches(9.3), Inches(2.5))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Place your quote text here. This template is designed for impactful single-quote slides."
    run.font.name = FONT_SERIF
    run.font.size = Pt(28)
    run.font.color.rgb = NAVY
    run.font.italic = True
    p.line_spacing = Pt(42)

    # Gold underline
    add_gold_accent_line(slide, Inches(2.0), Inches(5.2), Inches(2))

    # Attribution
    tf2 = slide.shapes.add_textbox(Inches(2.0), Inches(5.5), Inches(9), Inches(0.8))
    tf2.text_frame.word_wrap = True
    p2 = tf2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "— Attribution Name"
    run2.font.name = FONT_SANS
    run2.font.size = Pt(14)
    run2.font.color.rgb = GRAY_TEXT
    run2.font.bold = True

    p3 = tf2.text_frame.add_paragraph()
    run3 = p3.add_run()
    run3.text = "Title or Role"
    run3.font.name = FONT_MONO
    run3.font.size = Pt(10)
    run3.font.color.rgb = GOLD

    add_footer(slide, dark=False)


def slide_3_list(prs):
    """LIST SLIDE — Navy left panel, white right content area with bullet list"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_bg(slide)

    # Left navy panel
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()

    # Gold accent on panel
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.06), Inches(1.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # Section label on navy panel
    tf_label = slide.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(3), Inches(0.4))
    p_label = tf_label.text_frame.paragraphs[0]
    run_label = p_label.add_run()
    run_label.text = "SECTION"
    run_label.font.name = FONT_MONO
    run_label.font.size = Pt(10)
    run_label.font.color.rgb = GOLD

    # Title on navy panel
    tf_title = slide.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(3), Inches(1.5))
    tf_title.text_frame.word_wrap = True
    p_title = tf_title.text_frame.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "Slide Title Goes Here"
    run_title.font.name = FONT_SERIF
    run_title.font.size = Pt(32)
    run_title.font.color.rgb = WHITE

    # Right side — bullet list
    bullets = [
        ("First Point", "Supporting detail or explanation for this item"),
        ("Second Point", "Another supporting detail goes here"),
        ("Third Point", "Additional context or description"),
        ("Fourth Point", "Final supporting detail")
    ]

    y = Inches(1.6)
    for title, desc in bullets:
        # Gold bullet square
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), y + Inches(0.08), Inches(0.18), Inches(0.18))
        sq.fill.solid()
        sq.fill.fore_color.rgb = GOLD
        sq.line.fill.background()

        # Bullet title
        tf_b = slide.shapes.add_textbox(Inches(5.8), y - Inches(0.05), Inches(6.5), Inches(0.4))
        p_b = tf_b.text_frame.paragraphs[0]
        run_b = p_b.add_run()
        run_b.text = title
        run_b.font.name = FONT_SANS
        run_b.font.size = Pt(16)
        run_b.font.bold = True
        run_b.font.color.rgb = NAVY

        # Bullet description
        tf_d = slide.shapes.add_textbox(Inches(5.8), y + Inches(0.35), Inches(6.5), Inches(0.5))
        tf_d.text_frame.word_wrap = True
        p_d = tf_d.text_frame.paragraphs[0]
        run_d = p_d.add_run()
        run_d.text = desc
        run_d.font.name = FONT_SANS
        run_d.font.size = Pt(12)
        run_d.font.color.rgb = GRAY_TEXT

        y += Inches(1.2)

    add_footer(slide, dark=False)


def slide_4_image_left(prs):
    """IMAGE LEFT SLIDE — Image placeholder on left, text on right"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_bg(slide)

    # Top navy line
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Image placeholder (navy box)
    img_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.9), Inches(5.8), Inches(5.4))
    img_box.fill.solid()
    img_box.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xEE)
    img_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)
    img_box.line.width = Pt(1)

    # Placeholder text
    tf_ph = slide.shapes.add_textbox(Inches(2.0), Inches(3.2), Inches(3), Inches(0.6))
    p_ph = tf_ph.text_frame.paragraphs[0]
    p_ph.alignment = PP_ALIGN.CENTER
    run_ph = p_ph.add_run()
    run_ph.text = "Insert Image Here"
    run_ph.font.name = FONT_MONO
    run_ph.font.size = Pt(14)
    run_ph.font.color.rgb = RGBColor(0x99, 0x99, 0xAA)

    # Right side content
    add_corner_accent_at(slide, Inches(7.2), Inches(1.5))

    # Section label
    tf_label = slide.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5), Inches(0.3))
    p_l = tf_label.text_frame.paragraphs[0]
    run_l = p_l.add_run()
    run_l.text = "SECTION LABEL"
    run_l.font.name = FONT_MONO
    run_l.font.size = Pt(10)
    run_l.font.color.rgb = GOLD

    # Title
    tf_t = slide.shapes.add_textbox(Inches(7.2), Inches(2.5), Inches(5.3), Inches(1.2))
    tf_t.text_frame.word_wrap = True
    p_t = tf_t.text_frame.paragraphs[0]
    run_t = p_t.add_run()
    run_t.text = "Slide Title with Image"
    run_t.font.name = FONT_SERIF
    run_t.font.size = Pt(32)
    run_t.font.color.rgb = NAVY

    # Body text
    tf_body = slide.shapes.add_textbox(Inches(7.2), Inches(3.9), Inches(5.3), Inches(2))
    tf_body.text_frame.word_wrap = True
    p_body = tf_body.text_frame.paragraphs[0]
    run_body = p_body.add_run()
    run_body.text = "Add your description text here. This layout places a large image on the left with supporting text on the right. Ideal for screenshots, diagrams, or product photos."
    run_body.font.name = FONT_SANS
    run_body.font.size = Pt(14)
    run_body.font.color.rgb = GRAY_TEXT
    p_body.line_spacing = Pt(24)

    add_footer(slide, dark=False)


def slide_5_image_right(prs):
    """IMAGE RIGHT SLIDE — Text on left, image placeholder on right, navy accent"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_bg(slide)

    # Top gold line
    add_gold_accent_line(slide, 0, 0, W)

    # Left content
    # Gold accent square
    add_corner_accent_at(slide, Inches(0.8), Inches(1.3))

    # Section label
    tf_label = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5), Inches(0.3))
    p_l = tf_label.text_frame.paragraphs[0]
    run_l = p_l.add_run()
    run_l.text = "SECTION LABEL"
    run_l.font.name = FONT_MONO
    run_l.font.size = Pt(10)
    run_l.font.color.rgb = GOLD

    # Title
    tf_t = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(1.2))
    tf_t.text_frame.word_wrap = True
    p_t = tf_t.text_frame.paragraphs[0]
    run_t = p_t.add_run()
    run_t.text = "Another Layout with Image"
    run_t.font.name = FONT_SERIF
    run_t.font.size = Pt(32)
    run_t.font.color.rgb = NAVY

    # Body text
    tf_body = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(5.5), Inches(2.2))
    tf_body.text_frame.word_wrap = True
    p_body = tf_body.text_frame.paragraphs[0]
    run_body = p_body.add_run()
    run_body.text = "Add your description text here. This layout places the image on the right side. Use it to alternate visual flow between slides for a more dynamic presentation."
    run_body.font.name = FONT_SANS
    run_body.font.size = Pt(14)
    run_body.font.color.rgb = GRAY_TEXT
    p_body.line_spacing = Pt(24)

    # Gold accent line under text
    add_gold_accent_line(slide, Inches(0.8), Inches(5.8), Inches(2))

    # Right image placeholder
    img_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(0.9), Inches(5.8), Inches(5.4))
    img_box.fill.solid()
    img_box.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xEE)
    img_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)
    img_box.line.width = Pt(1)

    # Placeholder text
    tf_ph = slide.shapes.add_textbox(Inches(8.4), Inches(3.2), Inches(3), Inches(0.6))
    p_ph = tf_ph.text_frame.paragraphs[0]
    p_ph.alignment = PP_ALIGN.CENTER
    run_ph = p_ph.add_run()
    run_ph.text = "Insert Image Here"
    run_ph.font.name = FONT_MONO
    run_ph.font.size = Pt(14)
    run_ph.font.color.rgb = RGBColor(0x99, 0x99, 0xAA)

    add_footer(slide, dark=False)


def add_corner_accent_at(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.3), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def create_template():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_1_title(prs)
    slide_2_quote(prs)
    slide_3_list(prs)
    slide_4_image_left(prs)
    slide_5_image_right(prs)

    output = "/app/TrustOffice_Slide_Templates.pptx"
    prs.save(output)
    print(f"Saved to {output}")
    return output


if __name__ == "__main__":
    create_template()
